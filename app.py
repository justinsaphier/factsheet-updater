"""
SAM Factsheet Updater — Web App
================================
Run with:  python app.py
Then open: http://localhost:5000
"""

from flask import Flask, request, send_file, jsonify
import zipfile, io, shutil, os, tempfile
import numpy as np
import pandas as pd
import openpyxl
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50 MB


# ─────────────────────────────────────────────────────────────────────────────
# SECTOR MAP
# ─────────────────────────────────────────────────────────────────────────────
SECTOR_MAP = {
    "GOOGL": "Communications",
    "AVGO":  "Technology",
    "JPM":   "Financials",
    "CEG":   "Utilities",
    "MSFT":  "Technology",
    "LLY":   "Healthcare",
    "CSCO":  "Technology",
    "BRKB":  "Financials",
    "GD":    "Industrials",
    "LMT":   "Industrials",
    "HON":   "Industrials",
    "ICE":   "Financials",
    "JNJ":   "Healthcare",
    "EMR":   "Industrials",
    "ABBNY": "Industrials",
    "AON":   "Financials",
    "BHP":   "Materials",
    "CVX":   "Energy",
    "DIS":   "Communications",
    "KMI":   "Energy",
    "MCHP":  "Technology",
    "MELI":  "Cons. Discretionary",
    "NSRGY": "Consumer Staples",
    "PM":    "Consumer Staples",
}


# ─────────────────────────────────────────────────────────────────────────────
# DATE HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _to_short_date(long_date):
    month_map = {
        "January": "1", "February": "2", "March": "3",
        "April": "4", "May": "5", "June": "6",
        "July": "7", "August": "8", "September": "9",
        "October": "10", "November": "11", "December": "12"
    }
    for month, num in month_map.items():
        if long_date.startswith(month + " "):
            rest = long_date[len(month)+1:].replace(",", "").split()
            if len(rest) == 2:
                return f"{num}/{rest[0]}/{rest[1]}"
    return long_date

def _to_upper_date(long_date):
    parts = long_date.split(" ", 1)
    return parts[0].upper() + (" " + parts[1] if len(parts) > 1 else "")


# ─────────────────────────────────────────────────────────────────────────────
# PROCESSING
# ─────────────────────────────────────────────────────────────────────────────
def process_holdings(holdings_bytes, model_bytes):
    holdings = pd.read_excel(io.BytesIO(holdings_bytes))
    model    = pd.read_excel(io.BytesIO(model_bytes))

    model_tickers = set(model["Ticker"].astype(str).str.strip().str.upper())
    holdings["Ticker_clean"] = holdings["Ticker"].astype(str).str.strip().str.upper()

    filtered = holdings[
        holdings["Ticker_clean"].isin(model_tickers) &
        (holdings["Asset Category"] == "Equity")
    ].copy()

    if filtered.empty:
        raise ValueError(
            "No matching stocks found after filtering. "
            "Check that the holdings file has Asset Category and Ticker columns "
            "and the model file has a Ticker column."
        )

    stocks = filtered.groupby(["Ticker_clean", "ProductName"], as_index=False)["Total Value"].sum()
    total = stocks["Total Value"].sum()
    stocks["Weight_Pct"] = (stocks["Total Value"] / total * 100).round(2)
    top10 = stocks.sort_values("Weight_Pct", ascending=False).head(10).reset_index(drop=True)

    stocks["Sector"] = stocks["Ticker_clean"].map(SECTOR_MAP).fillna("Other")
    sector_weights = (
        stocks.groupby("Sector")["Weight_Pct"].sum()
        .reset_index().rename(columns={"Weight_Pct": "Sector_Pct"})
        .sort_values("Sector_Pct", ascending=False).reset_index(drop=True)
    )
    sector_weights["Sector_Pct"] = sector_weights["Sector_Pct"].round(1)
    return top10, sector_weights


def update_pie_chart_bytes(pptx_bytes, sector_weights_df):
    xlsx_key = "ppt/embeddings/Microsoft_Excel_Worksheet.xlsx"

    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
        names = zin.namelist()
        if xlsx_key not in names:
            raise ValueError("Could not find the embedded chart workbook in the PPTX.")
        xlsx_data = zin.read(xlsx_key)
        all_files = {name: zin.read(name) for name in names}

    wb = openpyxl.load_workbook(io.BytesIO(xlsx_data))
    if "Sectors" not in wb.sheetnames:
        raise ValueError(f"Embedded workbook has no 'Sectors' sheet. Found: {wb.sheetnames}")

    ws = wb["Sectors"]
    sectors = sector_weights_df["Sector"].tolist()
    pcts    = sector_weights_df["Sector_Pct"].tolist()

    for col in range(1, 20):
        ws.cell(row=1, column=col).value = None
        ws.cell(row=2, column=col).value = None
    for i, (sector, pct) in enumerate(zip(sectors, pcts)):
        ws.cell(row=1, column=i + 1).value = f"{sector}  {pct}%"
        ws.cell(row=2, column=i + 1).value = round(pct / 100, 4)

    buf = io.BytesIO()
    wb.save(buf)

    out_buf = io.BytesIO()
    with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in all_files.items():
            zout.writestr(name, buf.getvalue() if name == xlsx_key else data)

    pptx_bytes = out_buf.getvalue()

    chart_key = "ppt/charts/chart1.xml"
    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
        if chart_key not in zin.namelist():
            return pptx_bytes
        chart_xml = zin.read(chart_key)
        all_files = {name: zin.read(name) for name in zin.namelist()}

    root = etree.fromstring(chart_xml)
    for strCache in root.iter(qn("c:strCache")):
        for pt in strCache.findall(qn("c:pt")): strCache.remove(pt)
        ptc = strCache.find(qn("c:ptCount"))
        if ptc is not None: ptc.set("val", str(len(sectors)))
        for i, (s, p) in enumerate(zip(sectors, pcts)):
            pt = etree.SubElement(strCache, qn("c:pt"))
            pt.set("idx", str(i))
            etree.SubElement(pt, qn("c:v")).text = f"{s}  {p}%"

    for numCache in root.iter(qn("c:numCache")):
        for pt in numCache.findall(qn("c:pt")): numCache.remove(pt)
        ptc = numCache.find(qn("c:ptCount"))
        if ptc is not None: ptc.set("val", str(len(pcts)))
        for i, p in enumerate(pcts):
            pt = etree.SubElement(numCache, qn("c:pt"))
            pt.set("idx", str(i))
            etree.SubElement(pt, qn("c:v")).text = str(round(p / 100, 4))

    all_files[chart_key] = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    out_buf = io.BytesIO()
    with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in all_files.items(): zout.writestr(name, data)
    return out_buf.getvalue()


def _set_para_text(para, text):
    if para.runs:
        para.runs[0].text = str(text)
        for run in para.runs[1:]: run.text = ""
    else:
        para.text = str(text)


def replace_text(slide, old, new):
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    count += 1
    return count


def run_update(pptx_bytes, holdings_bytes, model_bytes, old_date, new_date):
    log = []

    log.append("Processing holdings...")
    top10, sector_weights = process_holdings(holdings_bytes, model_bytes)

    log.append("Top 10 holdings:")
    for _, row in top10.iterrows():
        log.append(f"  {row['ProductName']:<42} {row['Weight_Pct']:.2f}%")

    log.append("\nSector weights:")
    for _, row in sector_weights.iterrows():
        log.append(f"  {row['Sector']:<28} {row['Sector_Pct']:.1f}%")

    log.append("\nUpdating pie chart...")
    pptx_bytes = update_pie_chart_bytes(pptx_bytes, sector_weights)
    log.append("Pie chart updated.")

    log.append("\nUpdating Top 10 and dates...")
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]

    shapes_4 = sorted(
        [s for s in slide.shapes if s.name == "object 4" and s.has_text_frame],
        key=lambda s: s.left
    )
    if len(shapes_4) < 2:
        raise RuntimeError(f"Expected 2 Top 10 text boxes, found {len(shapes_4)}.")

    left_paras  = shapes_4[0].text_frame.paragraphs
    right_paras = shapes_4[1].text_frame.paragraphs

    for i in range(min(5, len(left_paras) - 1)):
        _set_para_text(left_paras[i + 1], top10.iloc[i]["ProductName"])
    top5_pct = int(round(top10.head(5)["Weight_Pct"].sum(), 0))
    _set_para_text(right_paras[0], f"MV {top5_pct}%")
    for i in range(min(5, len(right_paras) - 1)):
        _set_para_text(right_paras[i + 1], top10.iloc[5 + i]["ProductName"])
    log.append(f"Top 10 updated. MV of top 5 = {top5_pct}%")

    old_short = _to_short_date(old_date)
    new_short = _to_short_date(new_date)
    old_upper = _to_upper_date(old_date)
    new_upper = _to_upper_date(new_date)

    n1 = replace_text(slide, old_short, new_short)
    n2 = replace_text(slide, old_upper, new_upper)
    n3 = replace_text(slide, old_date,  new_date)
    log.append(f"Dates updated ({n1+n2+n3} replacements)")

    out_buf = io.BytesIO()
    prs.save(out_buf)

    return out_buf.getvalue(), log


# ─────────────────────────────────────────────────────────────────────────────
# ROUTES
# ─────────────────────────────────────────────────────────────────────────────
HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>SAM Factsheet Updater</title>
<style>
  *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
  :root {
    --bg:#F7F6F2; --surface:#FFFFFF; --border:#E2E0D8; --border2:#C8C6BC;
    --text:#1A1916; --muted:#6B6960; --accent:#1C5C3A; --accent2:#2E7D52;
    --accent-bg:#EBF4EF; --danger:#C0392B; --danger-bg:#FDF0EE;
    --warn:#7A5C00; --warn-bg:#FEF9EC; --radius:10px;
    --mono:"SF Mono","Fira Mono","Cascadia Code",monospace;
  }
  body { font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",system-ui,sans-serif;
    background:var(--bg);color:var(--text);min-height:100vh;display:flex;
    flex-direction:column;align-items:center;padding:40px 16px 80px; }
  .wordmark { font-size:11px;font-weight:700;letter-spacing:0.12em;color:var(--muted);
    text-transform:uppercase;margin-bottom:28px;display:flex;align-items:center;gap:10px; }
  .wordmark::before,.wordmark::after { content:"";display:block;height:1px;width:48px;background:var(--border2); }
  h1 { font-size:28px;font-weight:600;letter-spacing:-0.02em;text-align:center;margin-bottom:6px; }
  .subtitle { font-size:14px;color:var(--muted);text-align:center;margin-bottom:40px; }
  .card { background:var(--surface);border:1px solid var(--border);border-radius:var(--radius);
    padding:28px 32px;width:100%;max-width:580px;margin-bottom:16px; }
  .section-label { font-size:11px;font-weight:700;letter-spacing:0.1em;text-transform:uppercase;
    color:var(--muted);margin-bottom:20px;padding-bottom:10px;border-bottom:1px solid var(--border); }
  .field { margin-bottom:18px; }
  .field:last-child { margin-bottom:0; }
  .field label { display:block;font-size:13px;font-weight:600;margin-bottom:3px;color:var(--text); }
  .field .hint { font-size:12px;color:var(--muted);margin-bottom:7px;line-height:1.4; }
  .file-input-wrap { position:relative;display:flex;align-items:center; }
  .file-display { flex:1;height:38px;padding:0 12px;border:1px solid var(--border);border-right:none;
    border-radius:var(--radius) 0 0 var(--radius);font-size:13px;color:var(--muted);background:var(--bg);
    display:flex;align-items:center;overflow:hidden;white-space:nowrap;text-overflow:ellipsis;cursor:default; }
  .file-display.has-file { color:var(--text); }
  .file-input-wrap input[type="file"] { position:absolute;opacity:0;width:0;height:0; }
  .file-btn { height:38px;padding:0 14px;background:var(--surface);border:1px solid var(--border);
    border-radius:0 var(--radius) var(--radius) 0;font-size:13px;font-weight:500;color:var(--text);
    cursor:pointer;white-space:nowrap;transition:background 0.12s; }
  .file-btn:hover { background:var(--bg);border-color:var(--border2); }
  .date-row { display:grid;grid-template-columns:1fr 28px 1fr;align-items:end; }
  .date-row .arrow { display:flex;align-items:center;justify-content:center;height:38px;color:var(--muted);font-size:16px; }
  .date-row .field { margin-bottom:0; }
  input[type="text"] { width:100%;height:38px;padding:0 12px;border:1px solid var(--border);
    border-radius:var(--radius);font-size:13px;color:var(--text);background:var(--surface);
    outline:none;font-family:inherit; }
  input[type="text"]:focus { border-color:var(--accent); }
  .date-hint { font-size:11px;color:var(--muted);margin-top:7px;text-align:center; }
  .run-btn { width:100%;max-width:580px;height:48px;background:var(--accent);color:#fff;border:none;
    border-radius:var(--radius);font-size:15px;font-weight:600;cursor:pointer;display:flex;
    align-items:center;justify-content:center;gap:8px;margin-bottom:16px;transition:background 0.15s; }
  .run-btn:hover { background:var(--accent2); }
  .run-btn:disabled { opacity:0.5;cursor:not-allowed; }
  .spinner { width:18px;height:18px;border:2px solid rgba(255,255,255,0.35);border-top-color:#fff;
    border-radius:50%;animation:spin 0.7s linear infinite;display:none; }
  @keyframes spin { to { transform:rotate(360deg); } }
  .run-btn.loading .spinner { display:block; }
  .run-btn.loading .btn-text { display:none; }
  .result-card { width:100%;max-width:580px;border-radius:var(--radius);border:1px solid var(--border);
    overflow:hidden;display:none; }
  .result-card.visible { display:block; }
  .result-header { padding:14px 20px;display:flex;align-items:center;justify-content:space-between;gap:12px; }
  .result-header.success { background:var(--accent-bg);border-bottom:1px solid #C3E0CF; }
  .result-header.error { background:var(--danger-bg);border-bottom:1px solid #F5C6C2; }
  .result-status { display:flex;align-items:center;gap:8px;font-size:13px;font-weight:600; }
  .result-status.success { color:var(--accent); }
  .result-status.error { color:var(--danger); }
  .download-btn { height:32px;padding:0 14px;background:var(--accent);color:#fff;border:none;
    border-radius:6px;font-size:12px;font-weight:600;cursor:pointer;white-space:nowrap;
    display:flex;align-items:center;gap:6px;flex-shrink:0;transition:background 0.12s; }
  .download-btn:hover { background:var(--accent2); }
  .log-area { padding:16px 20px;background:#FAFAF8;font-family:var(--mono);font-size:11.5px;
    line-height:1.7;color:var(--muted);max-height:280px;overflow-y:auto;white-space:pre-wrap;
    border-top:1px solid var(--border); }
  .log-area .log-ok { color:var(--accent); }
  .log-area .log-err { color:var(--danger); }
  .log-area .log-warn { color:var(--warn); }
  .docs-note { width:100%;max-width:580px;background:var(--warn-bg);border:1px solid #E8D080;
    border-radius:var(--radius);padding:14px 18px;font-size:12.5px;color:var(--warn);line-height:1.6;margin-top:0; }
  .docs-note strong { font-weight:600; }
  svg.icon { width:14px;height:14px;flex-shrink:0; }
</style>
</head>
<body>
<div class="wordmark">Strategy Asset Managers</div>
<h1>Factsheet Updater</h1>
<p class="subtitle">Upload three files and enter the new quarter-end date — done in seconds.</p>

<div class="card">
  <div class="section-label">Files</div>
  <div class="field">
    <label>Factsheet</label>
    <div class="hint">The PowerPoint (.pptx) for the strategy you are updating</div>
    <div class="file-input-wrap">
      <div class="file-display" id="pptx-display">No file selected</div>
      <input type="file" id="pptx-input" accept=".pptx" onchange="setFile(this,'pptx')">
      <button class="file-btn" onclick="document.getElementById('pptx-input').click()">Browse</button>
    </div>
  </div>
  <div class="field">
    <label>Holdings</label>
    <div class="hint">The holdings export (.xlsx) for this strategy from your portfolio system</div>
    <div class="file-input-wrap">
      <div class="file-display" id="holdings-display">No file selected</div>
      <input type="file" id="holdings-input" accept=".xlsx" onchange="setFile(this,'holdings')">
      <button class="file-btn" onclick="document.getElementById('holdings-input').click()">Browse</button>
    </div>
  </div>
  <div class="field">
    <label>Model</label>
    <div class="hint">The model file (.xlsx) listing approved tickers for this strategy</div>
    <div class="file-input-wrap">
      <div class="file-display" id="model-display">No file selected</div>
      <input type="file" id="model-input" accept=".xlsx" onchange="setFile(this,'model')">
      <button class="file-btn" onclick="document.getElementById('model-input').click()">Browse</button>
    </div>
  </div>
</div>

<div class="card">
  <div class="section-label">Quarter-end date</div>
  <div class="date-row">
    <div class="field">
      <label>Current date in file</label>
      <input type="text" id="old-date" value="December 31, 2025" placeholder="December 31, 2025">
    </div>
    <div class="arrow">&#8594;</div>
    <div class="field">
      <label>New date</label>
      <input type="text" id="new-date" value="March 31, 2026" placeholder="March 31, 2026">
    </div>
  </div>
  <div class="date-hint">Format: Month DD, YYYY &mdash; e.g. December 31, 2025 &nbsp;|&nbsp; March 31, 2026 &nbsp;|&nbsp; June 30, 2026</div>
</div>

<button class="run-btn" id="run-btn" onclick="runUpdate()">
  <span class="spinner"></span>
  <span class="btn-text">Update Factsheet</span>
</button>

<div class="result-card" id="result-card">
  <div class="result-header" id="result-header">
    <div class="result-status" id="result-status"></div>
    <button class="download-btn" id="download-btn" style="display:none" onclick="downloadFile()">
      <svg class="icon" viewBox="0 0 16 16" fill="none" stroke="currentColor" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round">
        <path d="M8 2v8M5 7l3 3 3-3M2 12v1a1 1 0 001 1h10a1 1 0 001-1v-1"/>
      </svg>
      Download
    </button>
  </div>
  <div class="log-area" id="log-area"></div>
</div>

<div class="docs-note" style="margin-top:16px">
  <strong>Required file formats:</strong>
  Holdings must have columns: <code>Ticker</code>, <code>Asset Category</code>, <code>Total Value</code>, <code>ProductName</code>.
  Model must have a <code>Ticker</code> column.
  To run all five strategies, repeat with each strategy's own .pptx, holdings, and model files.
</div>

<script>
  let resultData = null;
  let resultFilename = null;

  function setFile(input, key) {
    const display = document.getElementById(key + '-display');
    if (input.files && input.files[0]) {
      display.textContent = input.files[0].name;
      display.classList.add('has-file');
    }
  }

  function escHtml(str) {
    return str.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
  }

  function formatLog(lines) {
    return lines.map(line => {
      if (line.includes('ERROR') || line.includes('error'))
        return '<span class="log-err">' + escHtml(line) + '</span>';
      if (line.includes('WARNING'))
        return '<span class="log-warn">' + escHtml(line) + '</span>';
      if (line.startsWith('DONE') || line.includes('updated') || line.includes('Pie chart'))
        return '<span class="log-ok">' + escHtml(line) + '</span>';
      return escHtml(line);
    }).join('\\n');
  }

  async function runUpdate() {
    const pptx     = document.getElementById('pptx-input').files[0];
    const holdings = document.getElementById('holdings-input').files[0];
    const model    = document.getElementById('model-input').files[0];
    const oldDate  = document.getElementById('old-date').value.trim();
    const newDate  = document.getElementById('new-date').value.trim();

    const missing = [];
    if (!pptx)     missing.push('Factsheet (.pptx)');
    if (!holdings) missing.push('Holdings (.xlsx)');
    if (!model)    missing.push('Model (.xlsx)');
    if (!oldDate)  missing.push('Current date');
    if (!newDate)  missing.push('New date');
    if (missing.length) { showError('Missing: ' + missing.join(', '), []); return; }

    const btn = document.getElementById('run-btn');
    btn.classList.add('loading');
    btn.disabled = true;

    const form = new FormData();
    form.append('pptx', pptx);
    form.append('holdings', holdings);
    form.append('model', model);
    form.append('old_date', oldDate);
    form.append('new_date', newDate);

    try {
      const res  = await fetch('/update', { method: 'POST', body: form });
      const json = await res.json();
      if (json.error) showError(json.error, json.log || []);
      else { resultData = json.data; resultFilename = json.filename; showSuccess(json.log || []); }
    } catch (e) {
      showError('Network error: ' + e.message, []);
    } finally {
      btn.classList.remove('loading');
      btn.disabled = false;
    }
  }

  function showSuccess(log) {
    const card = document.getElementById('result-card');
    card.classList.add('visible');
    document.getElementById('result-header').className = 'result-header success';
    document.getElementById('result-status').className = 'result-status success';
    document.getElementById('result-status').innerHTML =
      '<svg class="icon" viewBox="0 0 16 16" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M2 8l4 4 8-8"/></svg> Factsheet updated successfully';
    document.getElementById('download-btn').style.display = 'flex';
    document.getElementById('log-area').innerHTML = formatLog(log);
    card.scrollIntoView({ behavior: 'smooth', block: 'nearest' });
  }

  function showError(msg, log) {
    const card = document.getElementById('result-card');
    card.classList.add('visible');
    document.getElementById('result-header').className = 'result-header error';
    document.getElementById('result-status').className = 'result-status error';
    document.getElementById('result-status').innerHTML =
      '<svg class="icon" viewBox="0 0 16 16" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><circle cx="8" cy="8" r="6"/><path d="M8 5v3M8 11h.01"/></svg> ' + escHtml(msg);
    document.getElementById('download-btn').style.display = 'none';
    document.getElementById('log-area').innerHTML = formatLog(log);
    card.scrollIntoView({ behavior: 'smooth', block: 'nearest' });
  }

  function downloadFile() {
    if (!resultData) return;
    const bytes = new Uint8Array(resultData.match(/.{1,2}/g).map(b => parseInt(b, 16)));
    const blob  = new Blob([bytes], { type: 'application/vnd.openxmlformats-officedocument.presentationml.presentation' });
    const url   = URL.createObjectURL(blob);
    const a     = document.createElement('a');
    a.href = url;
    a.download = resultFilename || 'updated_factsheet.pptx';
    a.click();
    URL.revokeObjectURL(url);
  }
</script>
</body>
</html>"""


@app.route("/")
def index():
    return HTML


@app.route("/update", methods=["POST"])
def update():
    try:
        pptx_file     = request.files.get("pptx")
        holdings_file = request.files.get("holdings")
        model_file    = request.files.get("model")
        old_date      = request.form.get("old_date", "").strip()
        new_date      = request.form.get("new_date", "").strip()

        missing = []
        if not pptx_file:     missing.append("Factsheet (.pptx)")
        if not holdings_file: missing.append("Holdings (.xlsx)")
        if not model_file:    missing.append("Model (.xlsx)")
        if not old_date:      missing.append("Current date")
        if not new_date:      missing.append("New date")
        if missing:
            return jsonify({"error": f"Missing: {', '.join(missing)}"}), 400

        pptx_bytes     = pptx_file.read()
        holdings_bytes = holdings_file.read()
        model_bytes    = model_file.read()

        result_bytes, log = run_update(
            pptx_bytes, holdings_bytes, model_bytes, old_date, new_date
        )

        original_name = pptx_file.filename or "factsheet.pptx"
        base, ext = os.path.splitext(original_name)
        output_name = base.replace(old_date.split(",")[0].split(" ")[-1],
                                   new_date.split(",")[0].split(" ")[-1]) + ext
        if output_name == original_name:
            output_name = base + "_updated" + ext

        return jsonify({
            "log": log,
            "filename": output_name,
            "data": result_bytes.hex()
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    print(f"\n  SAM Factsheet Updater")
    print(f"  Open your browser to: http://localhost:{port}\n")
    app.run(debug=False, host="0.0.0.0", port=port)
